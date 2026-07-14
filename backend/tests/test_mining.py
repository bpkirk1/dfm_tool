"""Tests for the deterministic supplier-feedback miner (Phase 6)."""
from __future__ import annotations

import importlib.util
import json
from pathlib import Path

import pytest
import yaml

from app.mining import consolidate as C
from app.mining import emit as E
from app.mining import extract as X
from app.mining import model as M
from app.mining.model import Finding

MINING_DIR = Path(__file__).resolve().parents[2] / "new_suggestions" / "mining"

has_pptx = importlib.util.find_spec("pptx") is not None
has_xlsx = importlib.util.find_spec("openpyxl") is not None


# --- detectors ----------------------------------------------------------------
def test_detect_supplier_from_filename():
    assert M.detect_supplier("DFM for SentryHD RAM from AJET-2026-05-13.pptx") == "AJET"
    assert M.detect_supplier("PGCDDFMFor 757-4111-03X.pptx") == "Polygon"
    assert M.detect_supplier("HOKY 757-4311 A0.pptx") == "Hoky"
    assert M.detect_supplier("random.pptx") == "unknown"


def test_detect_date_variants():
    assert M.detect_date("deck-20260513.pptx") == "2026-05-13"
    assert M.detect_date("deck 2026-03-12.pptx") == "2026-03-12"
    assert M.detect_date("no date here.pptx") is None
    # a bare year must not be read as a date
    assert M.detect_date("thing_2026.pptx") is None


def test_detect_revision():
    assert M.detect_revision("x-sdcomments-x.pptx").lower().startswith("sdcomments")
    assert "rev" in M.detect_revision("Housing_Rev.04_.pptx").lower()


def test_detect_process_keywords():
    assert M.detect_process("increase inside bend radius, carrier tie bar") == "stamping"
    assert M.detect_process("gate vestige on molded LCP housing") == "molding"
    assert M.detect_process("selective plating PdNi zones") == "plating"
    assert M.detect_process("nothing relevant") == "unknown"


def test_detect_numbers_ignores_years():
    nums = M.detect_numbers("increase to 0.3mm, keep 45 deg, dated 2026")
    assert 0.3 in nums and 45.0 in nums
    assert 2026.0 not in nums


def test_parameter_hints():
    hits = M.detect_parameter_hints("the inside radius and burr height matter")
    assert "min_inside_corner_radius_mm" in hits
    assert "burr_height" in hits


# --- jsonl round-trip ---------------------------------------------------------
def test_jsonl_roundtrip(tmp_path):
    findings = [
        Finding(id="T-001", supplier="AJET", doc="d.pptx", slide=2, process="stamping",
                concern="x", parameter="burr_height", value_requested=0.02),
        Finding(id="T-002", supplier="AJET", doc="d.pptx", slide=3),
    ]
    p = tmp_path / "findings-ajet.jsonl"
    M.write_jsonl(findings, p)
    back = M.read_jsonl(p)
    assert [f.to_dict() for f in back] == [f.to_dict() for f in findings]


# --- consolidation (synthetic + real corpus) ---------------------------------
def _synthetic():
    return [
        Finding(id="A-1", supplier="AJET", doc="d1", process="stamping",
                parameter="min_inside_corner_radius_mm", value_requested=0.05,
                concern="radius"),
        Finding(id="P-1", supplier="Polygon", doc="d2", process="stamping",
                parameter="min_inside_corner_radius_mm", value_requested=0.08,
                value_accepted=0.08, resolution="accepted", part="757", date="2026-05-29",
                concern="radius accepted"),
        Finding(id="A-2", supplier="AJET", doc="d3", process="stamping",
                parameter="burr_height", value_requested=0.02, concern="burr"),
        Finding(id="A-3", supplier="AJET", doc="d4", confidence="image-only",
                concern="image"),
        Finding(id="A-4", supplier="AJET", doc="d5", process="assembly",
                parameter=None, concern="workflow note"),
    ]


def test_consolidation_grouping_conflict_and_capability():
    con = C.build_consolidation(_synthetic())
    radius = next(r for r in con.rule_candidates if r.parameter == "min_inside_corner_radius_mm")
    assert radius.seen_count == 2
    assert radius.conflict is True
    assert radius.value_min == 0.05 and radius.value_max == 0.08
    assert set(radius.suppliers) == {"AJET", "Polygon"}
    # capability picked up from the accepted value
    cap = next(c for c in con.capabilities if c.parameter == "min_inside_corner_radius_mm")
    assert cap.supplier == "Polygon" and cap.achieved == 0.08 and cap.resolution == "accepted"
    # routing
    assert any(i["id"] == "A-3" for i in con.image_only)
    assert any(w["id"] == "A-4" for w in con.workflow)


def test_consolidation_is_deterministic():
    a = C.build_consolidation(_synthetic()).to_dict()
    b = C.build_consolidation(_synthetic()).to_dict()
    assert a == b


@pytest.mark.skipif(not MINING_DIR.exists(), reason="no real mining corpus present")
def test_consolidation_over_real_findings():
    files = C.find_findings_files(MINING_DIR)
    assert files, "expected findings-*.jsonl in the corpus"
    con = C.build_consolidation(C.load_findings(files))
    assert con.rule_candidates
    md = C.render_markdown(con)
    assert md.startswith("# Stage 2")


# --- emit ---------------------------------------------------------------------
def test_emit_proposed_rules_are_all_proposed_and_parse():
    con = C.build_consolidation(_synthetic())
    text = E.emit_proposed_rules(con, run_date="2026-07-13")
    doc = yaml.safe_load(text)
    rules = [r for k, v in doc.items() if k.endswith("_rules") for r in v]
    assert rules
    assert all(r["status"] == "proposed" for r in rules)
    radius = next(r for r in rules if r["parameter"] == "min_inside_corner_radius_mm")
    assert radius["operator"] == "gte"          # inferred from "radius"
    assert radius["limit"] == 0.08              # most-conservative (max) for gte
    assert "conflict_note" in radius


def test_emit_operator_inference_ceiling():
    con = C.build_consolidation([
        Finding(id="A-1", supplier="AJET", doc="d", process="stamping",
                parameter="burr_height", value_requested=0.04, concern="b"),
        Finding(id="A-2", supplier="AJET", doc="d", process="stamping",
                parameter="burr_height", value_requested=0.02, concern="b"),
    ])
    rule = E.proposed_rule_dict(con.rule_candidates[0])
    assert rule["operator"] == "lte"
    assert rule["limit"] == 0.02  # smallest (most conservative) for a ceiling


def test_emit_ctf_and_review_queue_shapes():
    con = C.build_consolidation(_synthetic())
    ctf = json.loads(E.emit_ctf_entries(con))
    assert ctf["schema"] == "dfm-ctf-import/1"
    assert any(e["parameter"] == "min_inside_corner_radius_mm" for e in ctf["entries"])
    queue = E.emit_review_queue(con)
    assert "- [ ]" in queue and "## A. Proposed rules" in queue


# --- extraction ---------------------------------------------------------------
def test_extract_eml_uses_stdlib(tmp_path):
    eml = (
        "From: supplier@example.com\n"
        "Date: Mon, 18 Apr 2026 10:00:00 +0000\n"
        "Subject: DFM tie-bar vestige\n"
        "Content-Type: text/plain\n\n"
        "The sheared tie-bar vestige tolerance is +/-0.03mm and cannot be tighter.\n"
    )
    p = tmp_path / "AJET vestige 20260418.eml"
    p.write_text(eml, encoding="utf-8")
    findings, err, nxt = X.extract_file(p)
    assert err is None and findings
    f = findings[0]
    assert f.supplier == "AJET"
    assert f.date == "2026-04-18"
    assert "vestige" in f.concern.lower()
    assert "0.03" in f.note


def test_extract_unsupported_type_reports_error(tmp_path):
    p = tmp_path / "note.txt"
    p.write_text("hello", encoding="utf-8")
    findings, err, _ = X.extract_file(p)
    assert findings == [] and err and "unsupported" in err


@pytest.mark.skipif(not has_pptx, reason="python-pptx not installed")
def test_extract_pptx(tmp_path):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Increase inside radius to 0.08mm"
    p = tmp_path / "DFM from AJET 20260513.pptx"
    prs.save(str(p))
    findings, err, _ = X.extract_file(p)
    assert err is None
    assert any("radius" in f.concern.lower() for f in findings)


@pytest.mark.skipif(not has_xlsx, reason="openpyxl not installed")
def test_extract_xlsx(tmp_path):
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.append(["parameter", "value"])
    ws.append(["burr height", 0.02])
    p = tmp_path / "HOKY tracker 20260101.xlsx"
    wb.save(str(p))
    findings, err, _ = X.extract_file(p)
    assert err is None and findings
    assert findings[0].confidence == "table"

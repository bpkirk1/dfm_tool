"""Stage 1 — mechanical extraction of supplier documents into a findings scaffold.

Readers are lazy-imported so the core app has no new hard dependencies:

* .pptx  -> python-pptx (slide text, tables, speaker notes)
* .xlsx  -> openpyxl
* .eml   -> stdlib ``email``
* .pdf   -> pdfplumber (already used by the drawing extractor)

Each text unit (slide / sheet / email body / page) becomes one :class:`Finding`
with full provenance and detected hints. Semantic fields (parameter, resolution,
value) are left for the reviewer/LLM curation pass — the miner never guesses a
manufacturing value.
"""
from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Iterator

from . import model
from .model import Finding

_SUPPORTED = {".pptx", ".xlsx", ".eml", ".pdf"}


@dataclass
class TextUnit:
    """One extracted chunk of text with its location in a document."""
    slide: int | None
    kind: str          # "slide" | "table" | "notes" | "sheet" | "email" | "page"
    text: str
    image_only: bool = False


# --- per-format readers (lazy imports) ---------------------------------------
def _read_pptx(path: Path) -> tuple[list[TextUnit], str | None]:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as exc:
        return [], f"python-pptx unavailable: {exc}"
    units: list[TextUnit] = []
    prs = Presentation(str(path))
    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        has_picture = False
        for shape in slide.shapes:
            try:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    parts.append(shape.text_frame.text.strip())
                if getattr(shape, "has_table", False) and shape.has_table:
                    rows = []
                    for row in shape.table.rows:
                        rows.append(" | ".join(c.text.strip() for c in row.cells))
                    if rows:
                        units.append(TextUnit(i, "table", "\n".join(rows)))
                if shape.shape_type == 13:  # PICTURE
                    has_picture = True
            except Exception:
                continue
        body = "\n".join(parts).strip()
        if body:
            units.append(TextUnit(i, "slide", body))
        elif has_picture:
            units.append(TextUnit(i, "slide", "", image_only=True))
        # speaker notes
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    units.append(TextUnit(i, "notes", notes))
        except Exception:
            pass
    return units, None


def _read_xlsx(path: Path) -> tuple[list[TextUnit], str | None]:
    try:
        from openpyxl import load_workbook  # type: ignore
    except Exception as exc:
        return [], f"openpyxl unavailable: {exc}"
    units: list[TextUnit] = []
    wb = load_workbook(str(path), read_only=True, data_only=True)
    for ws in wb.worksheets:
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c not in (None, "")]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            units.append(TextUnit(None, "sheet", f"[{ws.title}]\n" + "\n".join(rows)))
    return units, None


def _read_eml(path: Path) -> tuple[list[TextUnit], str | None]:
    import email
    from email import policy

    msg = email.message_from_bytes(path.read_bytes(), policy=policy.default)
    header = " ".join(
        f"{h}: {msg.get(h)}" for h in ("From", "Date", "Subject") if msg.get(h)
    )
    body = ""
    try:
        part = msg.get_body(preferencelist=("plain", "html"))
        if part is not None:
            body = part.get_content()
    except Exception:
        body = msg.get_payload(decode=False) if isinstance(msg.get_payload(), str) else ""
    text = (header + "\n" + (body or "")).strip()
    return ([TextUnit(None, "email", text)] if text else []), None


def _read_pdf(path: Path) -> tuple[list[TextUnit], str | None]:
    try:
        import pdfplumber  # type: ignore
    except Exception as exc:
        return [], f"pdfplumber unavailable: {exc}"
    units: list[TextUnit] = []
    with pdfplumber.open(str(path)) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = (page.extract_text() or "").strip()
            if text:
                units.append(TextUnit(i, "page", text))
            else:
                units.append(TextUnit(i, "page", "", image_only=True))
    return units, None


_READERS = {
    ".pptx": _read_pptx,
    ".xlsx": _read_xlsx,
    ".eml": _read_eml,
    ".pdf": _read_pdf,
}


def read_document(path: str | Path) -> tuple[list[TextUnit], str | None]:
    """Return (text units, error). Error is a string if the file couldn't parse."""
    p = Path(path)
    reader = _READERS.get(p.suffix.lower())
    if reader is None:
        return [], f"unsupported file type: {p.suffix}"
    try:
        return reader(p)
    except Exception as exc:  # never abort a batch on one bad file
        return [], f"{type(exc).__name__}: {exc}"


# --- scaffold assembly --------------------------------------------------------
def _confidence(unit: TextUnit) -> str:
    if unit.image_only or not unit.text.strip():
        return "image-only"
    return "table" if unit.kind in ("table", "sheet") else "text"


def extract_file(
    path: str | Path, supplier: str | None = None, seq_start: int = 1
) -> tuple[list[Finding], str | None, int]:
    """Extract one document into findings. Returns (findings, error, next_seq)."""
    p = Path(path)
    units, err = read_document(p)
    if err is not None:
        return [], f"{p.name}: {err}", seq_start

    sup = supplier or model.detect_supplier(p.name)
    date = model.detect_date(p.name)
    revision = model.detect_revision(p.name)
    findings: list[Finding] = []
    seq = seq_start
    for unit in units:
        conf = _confidence(unit)
        text = unit.text.strip()
        hints = model.detect_parameter_hints(text) if text else []
        numbers = model.detect_numbers(text) if text else []
        process = model.detect_process(text) if text else "unknown"
        # A tight, provenance-preserving concern snippet (verbatim head of unit).
        concern = " ".join(text.split())[:280]
        note_bits = []
        if hints:
            note_bits.append("param_hints=" + ",".join(hints))
        if numbers:
            note_bits.append("numbers=" + ",".join(f"{n:g}" for n in numbers[:12]))
        note_bits.append(f"kind={unit.kind}")
        findings.append(
            Finding(
                id=f"{sup}-{seq:03d}",
                supplier=sup,
                doc=p.name,
                slide=unit.slide,
                date=date,
                revision=revision,
                process=process,
                concern=concern,
                parameter=None,
                resolution=model.guess_resolution(text) if text else "unknown",
                confidence=conf,
                note="; ".join(note_bits),
            )
        )
        seq += 1
    return findings, None, seq


def iter_documents(folder: str | Path) -> Iterator[Path]:
    for p in sorted(Path(folder).rglob("*")):
        if p.is_file() and p.suffix.lower() in _SUPPORTED and not p.name.startswith("~$"):
            yield p


def extract_folder(
    folder: str | Path, supplier: str | None = None
) -> tuple[list[Finding], list[str]]:
    """Extract every supported document under ``folder`` into findings.

    Returns (findings, unparsed) where ``unparsed`` lists documents that could
    not be read (missing optional reader, corrupt file, etc.).
    """
    findings: list[Finding] = []
    unparsed: list[str] = []
    # Per-supplier sequence counters keep ids stable and grouped.
    seq_by_supplier: dict[str, int] = {}
    for doc in iter_documents(folder):
        sup = supplier or model.detect_supplier(doc.name)
        start = seq_by_supplier.get(sup, 1)
        got, err, nxt = extract_file(doc, supplier=sup, seq_start=start)
        if err is not None:
            unparsed.append(err)
            continue
        findings.extend(got)
        seq_by_supplier[sup] = nxt
    return findings, unparsed
